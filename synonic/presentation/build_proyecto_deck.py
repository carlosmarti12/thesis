"""
Generator for synonic_proyecto_completo.pptx - full project overview deck:
objective, why this rebuild exists, the leakage fix, dataset, architecture,
the 11 methods, metrics, and results from the real 100-term comparison run
(results/comparison_results.csv). Reuses the same visual system as the other
decks in this folder (dark navy header bar, Calibri, same color tokens).

Run with: .venv/bin/python presentation/build_proyecto_deck.py
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x11, 0x1D, 0x3A)
BLUE = RGBColor(0x3D, 0x8B, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xC9, 0xD3, 0xE0)

SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)
MARGIN = Emu(502920)
CONTENT_TOP = Emu(1234440)
CONTENT_W = SLIDE_W - 2 * MARGIN

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def set_bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_header(slide, kicker, title, page_num):
    set_bg_white(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(1051560))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    kicker_box = slide.shapes.add_textbox(MARGIN, Emu(109728), Emu(10058400), Emu(320040))
    p = kicker_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = kicker
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(MARGIN, Emu(384048), Emu(10972800), Emu(640080))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = WHITE

    foot_l = slide.shapes.add_textbox(MARGIN, Emu(6510528), Emu(5486400), Emu(274320))
    p = foot_l.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "synonic — resumen del proyecto"
    r.font.size = Pt(10)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY

    foot_r = slide.shapes.add_textbox(Emu(11155680), Emu(6510528), Emu(548640), Emu(274320))
    p = foot_r.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(page_num)
    r.font.size = Pt(10)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY


def add_bullets(slide, items, top=Emu(1280160), height=Emu(5120640), left=Emu(640080), width=Emu(10881360), lvl0_size=17, lvl1_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for level, text, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        prefix = "•  " if level == 0 else "‒  "
        r = p.add_run()
        r.text = prefix + text
        r.font.name = "Calibri"
        r.font.bold = bold
        r.font.color.rgb = INK
        if level == 0:
            r.font.size = Pt(lvl0_size)
            p.space_after = Pt(8)
        else:
            r.font.size = Pt(lvl1_size)
            p.space_after = Pt(4)


def add_table(slide, headers, rows, col_widths, top=CONTENT_TOP, left=MARGIN, height=Emu(4846320), font_pt=13, header_font_pt=14, center_cols=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = sum(col_widths)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    center_cols = center_cols or set()

    for c, w in enumerate(col_widths):
        table.columns[c].width = w

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = htext
        r.font.size = Pt(header_font_pt)
        r.font.bold = True
        r.font.name = "Calibri"
        r.font.color.rgb = WHITE

    for ridx, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ridx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if c in center_cols:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = Pt(font_pt)
            r.font.bold = False
            r.font.name = "Calibri"
            r.font.color.rgb = INK


def title_slide(kicker_text, main_title, subtitle, footnote):
    slide = new_slide()
    set_bg_white(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = slide.shapes.add_textbox(MARGIN, Emu(2100000), Emu(11185855), Emu(700000))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = main_title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = WHITE

    box = slide.shapes.add_textbox(MARGIN, Emu(2870000), Emu(11185855), Emu(500000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(19)
    r.font.name = "Calibri"
    r.font.color.rgb = BLUE

    box = slide.shapes.add_textbox(MARGIN, Emu(3420000), Emu(11185855), Emu(500000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = footnote
    r.font.size = Pt(14)
    r.font.name = "Calibri"
    r.font.color.rgb = LIGHT

    box = slide.shapes.add_textbox(MARGIN, Emu(6300000), Emu(9000000), Emu(400000))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Carlos Martínez  ·  Proyecto synonic/"
    r.font.size = Pt(12)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY
    return slide


# ===========================================================================
# Load real results
# ===========================================================================
df = pd.read_csv("results/comparison_results.csv")
agg = df.groupby("method").agg(
    exact_top1=("exact_top1", "mean"),
    exact_any=("exact_any", "mean"),
    recall_at_5=("recall_at_5", "mean"),
    mrr=("mrr", "mean"),
    avg_time_s=("time_seconds", "mean"),
    avg_candidates=("num_final_candidates", "mean"),
)
ORDER = ["same_term", "wordnet_direct", "embedding_wordnet", "llm_zero_shot",
          "llm_expansion", "llm_rerank", "hybrid_fusion", "mas_base",
          "mas_llm_ranker", "mas_safe_hybrid", "supervisor"]
agg = agg.loc[ORDER]

# ===========================================================================
# Slide 1 — Title
# ===========================================================================
title_slide(
    "SYNONIC",
    "SYNONIC — RESUMEN DEL PROYECTO",
    "Objetivo, arquitectura, métodos y métricas de una tarea de descubrimiento abierto de sinónimos",
    "Qué hemos hecho, cómo lo hemos hecho y cómo lo medimos",
)

# ===========================================================================
# Slide 2 — Agenda
# ===========================================================================
slide = new_slide()
add_header(slide, "ÍNDICE", "Agenda", 2)
add_bullets(slide, [
    (0, "Objetivo y definición de la tarea", True),
    (0, "Por qué existe este rebuild — el problema de data leakage", True),
    (0, "Dataset y muestra de evaluación", True),
    (0, "Arquitectura general del proyecto", True),
    (0, "Los 11 métodos implementados", True),
    (0, "Cómo medimos — las métricas", True),
    (0, "Resultados (100 términos, llama3.2:3b)", True),
    (0, "Hallazgos clave y limitaciones", True),
    (0, "Próximos pasos", True),
])

# ===========================================================================
# Slide 3 — Objetivo
# ===========================================================================
slide = new_slide()
add_header(slide, "OBJETIVO", "Definición de la tarea", 3)
add_bullets(slide, [
    (0, "Open synonym discovery", True),
    (1, "Dado un término académico (p.ej. \"stocks\"), encontrar sus sinónimos o etiquetas alternativas tal y como aparecerían en una taxonomía/base de datos académica", False),
    (1, "La referencia (ground truth) es la columna en_synonym del dataset TU-Expert-Collection-Topic-Synonyms", False),
    (0, "Regla central: descubrimiento abierto", True),
    (1, "Un método solo puede ver el término de entrada (en) — nunca en_synonym hasta el momento de evaluar", False),
    (1, "en_synonym únicamente sirve para puntuar un resultado ya producido, nunca para ayudar a producirlo", False),
    (0, "Por qué importa", True),
    (1, "Simula un despliegue real: en producción no existe una \"respuesta correcta\" a la que buscar dentro", False),
])

# ===========================================================================
# Slide 4 — Por qué existe este rebuild (leakage)
# ===========================================================================
slide = new_slide()
add_header(slide, "RIGOR METODOLÓGICO", "El problema de data leakage y su corrección", 4)
add_bullets(slide, [
    (0, "Hallazgo en la iteración anterior (new/src2/)", True),
    (1, "Varios métodos construían su índice de búsqueda directamente sobre df['en_synonym'] — la columna de respuestas", False),
    (1, "Afectaba a: embedding retrieval, expansión LLM, y las 3 variantes MAS (su retrieval_agent)", False),
    (1, "Resultado: el sistema buscaba literalmente dentro de la respuesta correcta → métricas infladas, no representativas de un uso real", False),
    (0, "Corrección aplicada en synonic/", True),
    (1, "El vocabulario de candidatos ahora viene de WordNet (src/open_vocab.py) — recurso totalmente independiente del dataset", False),
    (1, "en_synonym se lee en un único punto de todo el código: run_experiment.py, dentro del bucle principal, siempre después de generar los candidatos", False),
    (1, "Verificable en cualquier momento: grep en_synonym src/ — cualquier otra aparición es solo comentario/docstring", False),
])

# ===========================================================================
# Slide 5 — Dataset
# ===========================================================================
slide = new_slide()
add_header(slide, "DATOS", "Dataset y muestra de evaluación", 5)
add_bullets(slide, [
    (0, "Fuente", True),
    (1, "jensjorisdecorte/TU-Expert-Collection-Topic-Synonyms (HuggingFace)", False),
    (1, "Se conservan solo topic, en (término de entrada) y en_synonym (referencia) — se descartan nl/nl_synonym (el proyecto es solo en inglés)", False),
    (0, "Preparación", True),
    (1, "data/prepare_dataset.py → data/topic_synonyms_clean.csv — 970 filas tras eliminar vacíos y duplicados", False),
    (1, "data/make_eval_sample.py → data/eval_sample.csv — muestra fija de 100 filas, semilla fija (random_state=42)", False),
    (1, "La semilla fija garantiza que todos los métodos se comparan sobre exactamente las mismas filas, de forma reproducible", False),
])

# ===========================================================================
# Slide 6 — Arquitectura general
# ===========================================================================
slide = new_slide()
add_header(slide, "ARQUITECTURA", "Cómo está organizado el proyecto", 6)
add_bullets(slide, [
    (0, "Flujo general", True),
    (1, "data/ → dataset limpio + muestra fija de evaluación", False),
    (1, "src/open_vocab.py → índice de vocabulario WordNet (~147k frases), cacheado en data/cache/", False),
    (1, "src/llm.py → único punto de acceso a Ollama (ask_llm, parse_candidates)", False),
    (1, "src/baselines.py + src/agents.py/graph.py + src/supervisor/ → los 11 métodos", False),
    (1, "src/evaluation.py → todas las métricas, documentadas", False),
    (1, "src/run_experiment.py (un método) / src/run_comparison.py (los 11 a la vez) → resultados/*.csv", False),
    (0, "Diseño clave", True),
    (1, "Todo método recibe solo term + recursos compartidos (embedder/LLM/índice) — nunca el dataframe completo", False),
    (1, "Ambos CLIs son resumibles: si el CSV ya tiene filas para (término, método), se saltan — seguro interrumpir y continuar", False),
])

# ===========================================================================
# Slide 7 — Los 11 métodos (tabla resumen)
# ===========================================================================
slide = new_slide()
add_header(slide, "MÉTODOS", "Los 11 métodos implementados (1/2)", 7)
add_table(
    slide,
    ["Método", "Descripción"],
    [
        ("same_term", "Devuelve el propio término. Suelo (floor) de referencia."),
        ("wordnet_direct", "Lemas de los synsets de WordNet que contienen el término. Simbólico, sin ML."),
        ("embedding_wordnet", "Embedding del término (MiniLM) + top-5 vecinos por coseno en el índice WordNet."),
        ("llm_zero_shot", "Pide 5 sinónimos directamente al LLM, sin recuperación."),
        ("llm_expansion", "El LLM genera 6 variantes de consulta; se buscan en WordNet y se agregan por score."),
        ("llm_rerank", "El LLM genera 10 candidatos libres; se reordenan por similitud de embedding."),
        ("hybrid_fusion", "3 llamadas LLM (self-consistency) + recuperación WordNet, fusionadas con RRF ponderado."),
    ],
    [Emu(2800000), Emu(8355680)],
    height=Emu(4600000),
    font_pt=13,
)

# ===========================================================================
# Slide 8 — Los 11 métodos (MAS + supervisor)
# ===========================================================================
slide = new_slide()
add_header(slide, "MÉTODOS", "Los 11 métodos implementados (2/2) — MAS y supervisor", 8)
add_bullets(slide, [
    (0, "mas_base (v1)", True),
    (1, "generator(LLM) + retrieval(WordNet) → merge → verification (0.75·semántico + 0.25·fuzzy) → top 5", False),
    (0, "mas_llm_ranker (v2)", True),
    (1, "retrieval → generator → merge → un LLM reordena la lista completa → top 5. Ablación negativa confirmada", False),
    (0, "mas_safe_hybrid (v3)", True),
    (1, "Igual que v2 + safe_finalizer: 0.9·embedding_sim + 0.1·bonus del LLM (nunca penaliza)", False),
    (0, "supervisor", True),
    (1, "Un único LLM decide, turno a turno, qué acción tomar (retrieve/generate/merge/score/finalize/finish)", False),
    (1, "Reutiliza las mismas funciones deterministas que los MAS — solo cambia quién decide el orden", False),
    (1, "Bug de planificación real encontrado y corregido con el propio harness de comparación (ver hallazgos)", False),
])

# ===========================================================================
# Slide 9 — Métricas
# ===========================================================================
slide = new_slide()
add_header(slide, "CÓMO MEDIMOS", "Definición de las métricas", 9)
add_table(
    slide,
    ["Métrica", "Qué mide"],
    [
        ("exact_top1", "El candidato #1 coincide exactamente con la referencia."),
        ("exact_any / recall@5", "La referencia aparece en cualquier posición de los 5 candidatos devueltos."),
        ("recall_at_3", "La referencia aparece entre los 3 primeros candidatos."),
        ("precision_at_5", "Solo hay un sinónimo de referencia por término → tope mecánico de 0.2; equivale a recall@5 / 5."),
        ("mrr", "1 / posición de la respuesta correcta (0 si no aparece). Premia acertar cerca del top aunque no sea #1."),
        ("fuzzy_similarity", "rapidfuzz.fuzz.ratio — cercanía léxica/ortográfica (p.ej. \"optimise\" vs \"optimize\")."),
        ("semantic_similarity", "Coseno de embeddings (MiniLM) — cercanía de significado, no de forma textual."),
    ],
    [Emu(2900000), Emu(8255680)],
    height=Emu(4600000),
    font_pt=13,
)

# ===========================================================================
# Slide 10 — Resultados (tabla completa, datos reales)
# ===========================================================================
slide = new_slide()
add_header(slide, "RESULTADOS", "Comparativa completa — 100 términos, llama3.2:3b", 10)


def fmt(x, pct=False):
    return f"{x:.2f}" if not pct else f"{x*100:.0f}%"


rows = []
for m in ORDER:
    row = agg.loc[m]
    rows.append((
        m,
        f"{row['mrr']:.3f}",
        f"{row['exact_any']*100:.0f}%",
        f"{row['recall_at_5']*100:.0f}%",
        f"{row['avg_time_s']:.2f}s",
    ))

add_table(
    slide,
    ["Método", "MRR", "exact_any", "recall@5", "Tiempo/término"],
    rows,
    [Emu(3200000), Emu(1996464), Emu(1996464), Emu(1996464), Emu(1996464)],
    height=Emu(4600000),
    font_pt=13,
    header_font_pt=13,
    center_cols={1, 2, 3, 4},
)

# ===========================================================================
# Slide 11 — Hallazgos clave
# ===========================================================================
slide = new_slide()
add_header(slide, "HALLAZGOS", "Lectura de los resultados", 11)
add_bullets(slide, [
    (0, "El patrón se repite en todas las iteraciones de la tesis", True),
    (1, "Los métodos de una sola llamada al LLM (llm_zero_shot, hybrid_fusion, llm_rerank) superan en MRR a las arquitecturas multiagente completas", False),
    (1, "Más \"autonomía\" agéntica (MAS, supervisor) no mejora los resultados frente a pipelines más simples — y cuesta más tiempo", False),
    (0, "hybrid_fusion es el mejor en recall@5 / exact_any (31%)", True),
    (1, "Combinar generación LLM con self-consistency + recuperación WordNet vía RRF gana en cobertura, aunque no en MRR puro", False),
    (0, "mas_llm_ranker (v2) confirma la ablación negativa", True),
    (1, "Dejar que el LLM reordene libremente la lista fusionada empeora el resultado frente a un ranking numérico (v1) o híbrido (v3)", False),
    (0, "supervisor: bug de planificación real, encontrado y corregido", True),
    (1, "El LLM a veces generaba candidatos después de fusionar (\"merge\"), y esos candidatos quedaban huérfanos sin re-fusionar", False),
    (1, "Tras el fix (comprobar orden, no solo presencia, de acciones): mrr 0.091 → 0.122, recall@5 0.13 → 0.21", False),
    (0, "Coste temporal", True),
    (1, "De 0s (same_term) a 7.2s/término (supervisor) — el método más lento no es el que mejor rinde", False),
])

# ===========================================================================
# Slide 12 — Limitaciones
# ===========================================================================
slide = new_slide()
add_header(slide, "LIMITACIONES", "Qué no cubre (todavía) este análisis", 12)
add_bullets(slide, [
    (0, "Modelo pequeño", True),
    (1, "llama3.2:3b vía Ollama — un modelo más grande probablemente mejoraría todos los métodos dependientes de LLM sin cambiar la arquitectura", False),
    (0, "Techo de cobertura de WordNet", True),
    (1, "No contiene todos los términos específicos de dominio que espera el dataset — penaliza a los métodos anclados al índice (retrieval, MAS) frente a los puramente generativos", False),
    (0, "Referencias idiosincráticas", True),
    (1, "Algunas respuestas de referencia son elecciones editoriales muy específicas que ningún método puede predecir solo a partir del término", False),
    (0, "Escala de la muestra", True),
    (1, "100 de 970 términos disponibles — validar a mayor escala es el siguiente paso natural", False),
])

# ===========================================================================
# Slide 13 — Próximos pasos
# ===========================================================================
slide = new_slide()
add_header(slide, "CIERRE", "Próximos pasos", 13)
add_bullets(slide, [
    (0, "Validar a mayor escala", True),
    (1, "Repetir la comparación completa sobre las 970 filas del dataset, no solo la muestra de 100", False),
    (0, "Probar un modelo más grande", True),
    (1, "Medir cuánto del gap entre métodos se debe al tamaño del LLM y no a la arquitectura", False),
    (0, "Profundizar en supervisor", True),
    (1, "Evaluar si separar el supervisor en agentes cognitivos independientes (desambiguación / generación / crítica) aporta valor real antes de construirlo", False),
    (0, "Documentar para la memoria de tesis", True),
    (1, "Trasladar esta comparación y sus hallazgos al capítulo de resultados, con las tablas completas de comparison_results.csv", False),
])

# ===========================================================================
# Slide 14 — Cierre
# ===========================================================================
slide = new_slide()
set_bg_white(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()
bar.shadow.inherit = False

box = slide.shapes.add_textbox(MARGIN, Emu(2900000), Emu(11185855), Emu(700000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Preguntas"
r.font.size = Pt(40)
r.font.bold = True
r.font.name = "Calibri"
r.font.color.rgb = WHITE

box = slide.shapes.add_textbox(MARGIN, Emu(3700000), Emu(11185855), Emu(500000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Código y resultados: proyecto synonic/"
r.font.size = Pt(16)
r.font.name = "Calibri"
r.font.color.rgb = BLUE

out_path = "presentation/synonic_proyecto_completo.pptx"
prs.save(out_path)
print("saved:", out_path)
print(agg.to_string())
