"""
One-off generator for synonic_metodos.pptx - reuses the same visual system as
synonic_presentation.pptx (dark navy header bar, Calibri, same color tokens)
but is a standalone deck focused only on explaining the 11 methods and how
they differ. Not wired into any pipeline; run once with:
    .venv/bin/python presentation/build_metodos_deck.py
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x11, 0x1D, 0x3A)
BLUE = RGBColor(0x3D, 0x8B, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x5B, 0x66)

SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)
MARGIN = Emu(502920)
CONTENT_TOP = Emu(1234440)
CONTENT_W = SLIDE_W - 2 * MARGIN

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def set_bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_header(slide, kicker, title, page_num):
    set_bg_white(slide)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(1051560))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    kicker_box = slide.shapes.add_textbox(MARGIN, Emu(109728), Emu(10058400), Emu(320040))
    tf = kicker_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = kicker
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(MARGIN, Emu(384048), Emu(10972800), Emu(640080))
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = WHITE

    foot_l = slide.shapes.add_textbox(MARGIN, Emu(6510528), Emu(5486400), Emu(274320))
    p = foot_l.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "synonic — explicación de métodos"
    r.font.size = Pt(10)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY

    foot_r = slide.shapes.add_textbox(Emu(11155680), Emu(6510528), Emu(548640), Emu(274320))
    p = foot_r.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(page_num)
    r.font.size = Pt(10)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY


def add_bullets(slide, items, top=Emu(1280160), height=Emu(5120640), left=Emu(640080), width=Emu(10881360)):
    """items: list of (level, text, bold) tuples. level 0 -> '•  ', level 1 -> '‒  '."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for level, text, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        prefix = "•  " if level == 0 else "‒  "
        r = p.add_run()
        r.text = prefix + text
        r.font.name = "Calibri"
        r.font.bold = bold
        r.font.color.rgb = INK
        if level == 0:
            r.font.size = Pt(17)
            p.space_after = Pt(8)
        else:
            r.font.size = Pt(15)
            p.space_after = Pt(4)


def add_table(slide, headers, rows, col_widths, top=CONTENT_TOP, left=MARGIN, height=Emu(4846320)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = sum(col_widths)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    for c, w in enumerate(col_widths):
        table.columns[c].width = w

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = htext
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.name = "Calibri"
        r.font.color.rgb = WHITE

    for ridx, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ridx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(13)
            r.font.bold = False
            r.font.name = "Calibri"
            r.font.color.rgb = INK


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
slide = new_slide()
set_bg_white(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()
bar.shadow.inherit = False

box = slide.shapes.add_textbox(MARGIN, Emu(2438400), Emu(11185855), Emu(700000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "MÉTODOS DE SYNONIC"
r.font.size = Pt(40)
r.font.bold = True
r.font.name = "Calibri"
r.font.color.rgb = WHITE

box = slide.shapes.add_textbox(MARGIN, Emu(3200400), Emu(11185855), Emu(500000))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Qué hace cada uno de los 11 métodos implementados y en qué se diferencian"
r.font.size = Pt(20)
r.font.name = "Calibri"
r.font.color.rgb = BLUE

box = slide.shapes.add_textbox(MARGIN, Emu(3750000), Emu(11185855), Emu(500000))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Baselines, recuperación por embeddings, LLM, arquitecturas MAS de grafo fijo y supervisor dinámico"
r.font.size = Pt(14)
r.font.name = "Calibri"
r.font.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)

box = slide.shapes.add_textbox(MARGIN, Emu(6300000), Emu(9000000), Emu(400000))
tf = box.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Carlos Martínez  ·  Proyecto synonic/"
r.font.size = Pt(12)
r.font.name = "Calibri"
r.font.color.rgb = GRAY

# ---------------------------------------------------------------------------
# Slide 2 — Agenda
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "ÍNDICE", "Agenda", 2)
add_bullets(slide, [
    (0, "Baselines sin ML — same_term, wordnet_direct", True),
    (0, "Recuperación por embeddings — embedding_wordnet", True),
    (0, "Un único LLM — llm_zero_shot, llm_rerank", True),
    (0, "Expansión de consulta y fusión — llm_expansion, hybrid_fusion", True),
    (0, "MAS de grafo fijo — mas_base, mas_llm_ranker, mas_safe_hybrid", True),
    (0, "Supervisor dinámico — supervisor", True),
    (0, "Comparativa por ejes clave y resultados", True),
])

# ---------------------------------------------------------------------------
# Slide 3 — Baselines sin ML
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "MÉTODOS · FAMILIA 1", "Baselines sin ML (src/baselines.py)", 3)
add_table(
    slide,
    ["Método", "Descripción técnica"],
    [
        ("same_term", "Devuelve el propio término tal cual. Suelo (floor): cualquier método real debe superarlo."),
        ("wordnet_direct", "Busca el término en WordNet y devuelve los lemas de los synsets donde aparece. Sin embeddings, sin LLM, sin vocabulario del dataset."),
    ],
    [Emu(2800000), Emu(8355680)],
    height=Emu(2200000),
)
add_bullets(slide, [
    (1, "Sirven como referencia mínima: cualquier método con ML/LLM debería superarlos claramente", False),
], top=Emu(3650000), height=Emu(1500000))

# ---------------------------------------------------------------------------
# Slide 4 — Embeddings + un único LLM
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "MÉTODOS · FAMILIA 2 Y 3", "Embeddings puros y un único LLM", 4)
add_table(
    slide,
    ["Método", "Descripción técnica"],
    [
        ("embedding_wordnet", "Embebe el término (MiniLM) y busca los top-5 vecinos por coseno en un índice precalculado de todo el vocabulario de WordNet (~147k frases). Línea base 'embeddings puros'."),
        ("llm_zero_shot", "Pide directamente al LLM 5 sinónimos, sin recuperación de ningún tipo. Confía en el orden que da el propio modelo."),
        ("llm_rerank", "El LLM genera 10 candidatos libres (no restringidos a WordNet); se reordenan numéricamente por similitud de embedding con el término. No confía en el orden del LLM."),
    ],
    [Emu(2800000), Emu(8355680)],
    height=Emu(3200000),
)

# ---------------------------------------------------------------------------
# Slide 5 — Expansión de consulta y fusión
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "MÉTODOS · FAMILIA 4", "Expansión de consulta + fusión", 5)
add_bullets(slide, [
    (0, "llm_expansion", True),
    (1, "El LLM genera 6 variantes/paráfrasis del término (consultas, no candidatos finales)", False),
    (1, "Cada variante se embebe y se busca en el índice de WordNet", False),
    (1, "Se agrega por score máximo entre variantes + bonus si un candidato aparece en varias", False),
    (1, "Diferencia con llm_rerank: el LLM nunca propone el sinónimo final, solo amplía la consulta", False),
    (0, "hybrid_fusion", True),
    (1, "El más elaborado de los no-agénticos: 3 llamadas al LLM con temperatura>0 (self-consistency) → lista A", False),
    (1, "+ recuperación por embeddings en WordNet → lista B", False),
    (1, "Fusiona A y B con Reciprocal Rank Fusion ponderado (peso 1.0 al LLM, 0.5 a WordNet)", False),
    (1, "Combina lo mejor de llm_rerank (candidatos libres) y embedding_wordnet (candidatos garantizados)", False),
])

# ---------------------------------------------------------------------------
# Slide 6 — MAS de grafo fijo
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "MÉTODOS · FAMILIA 5", "MAS de grafo fijo (LangGraph)", 6)
add_bullets(slide, [
    (0, "mas_base (v1)", True),
    (1, "generator(LLM) → retrieval(WordNet) → merge → verification (0.75·semantic + 0.25·fuzzy) → ranking → top-5", False),
    (1, "Puntuación final numérica; el LLM solo genera candidatos, nunca ordena", False),
    (0, "mas_llm_ranker (v2)", True),
    (1, "retrieval → generator → merge (cap 20) → un LLM reordena la lista completa → top-5", False),
    (1, "Ablación negativa confirmada: confiar el ranking final al LLM empeora el resultado", False),
    (0, "mas_safe_hybrid (v3)", True),
    (1, "Igual que v2 + safe_finalizer: score = 0.9·embedding_sim + 0.1·llm_rank_bonus", False),
    (1, "El LLM solo puede sumar bonus — nunca penaliza un candidato que dejó fuera del reordenamiento", False),
])

# ---------------------------------------------------------------------------
# Slide 7 — Supervisor dinámico
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "MÉTODOS · FAMILIA 6", "Supervisor dinámico (src/supervisor/)", 7)
add_bullets(slide, [
    (0, "supervisor", True),
    (1, "Una única llamada al LLM por turno decide la siguiente acción: retrieve / generate / merge / score / finalize / finish", False),
    (1, "Reutiliza las mismas funciones de agents.py como herramientas — no reimplementa lógica", False),
    (1, "Control de flujo dinámico frente al orden fijo de nodos de las variantes MAS", False),
    (0, "Bug de planificación encontrado y corregido", True),
    (1, "El modelo a veces elegía 'generate' después de 'merge': esos candidatos quedaban huérfanos sin re-fusionar", False),
    (1, "Fix: los guardrails ahora comprueban el ORDEN de las acciones, no solo si ya se ejecutaron", False),
    (1, "Tras el fix: mrr 0.091 → 0.1215 (+34%), recall@5 0.13 → 0.21 — pasa de peor de 11 a empate con los demás MAS", False),
])

# ---------------------------------------------------------------------------
# Slide 8 — Comparativa por ejes clave
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "COMPARATIVA", "Diferencias clave entre familias", 8)
add_table(
    slide,
    ["Eje", "Métodos"],
    [
        ("¿Usa LLM?", "No: same_term, wordnet_direct, embedding_wordnet.  Sí: el resto (8 métodos)."),
        ("¿Restringido a vocabulario WordNet?", "Sí: wordnet_direct, embedding_wordnet, llm_expansion, los 3 MAS, supervisor.  No (LLM libre): llm_zero_shot, llm_rerank, parte de hybrid_fusion."),
        ("¿Quién decide el ranking final?", "Numérico/embeddings: mas_base, mas_safe_hybrid, llm_rerank, embedding_wordnet.  El LLM directamente: llm_zero_shot, mas_llm_ranker."),
        ("¿Orden de pasos fijo o dinámico?", "Fijo: todos excepto supervisor.  Dinámico (decidido por el LLM turno a turno): supervisor."),
    ],
    [Emu(3800000), Emu(7355680)],
    height=Emu(3800000),
)

# ---------------------------------------------------------------------------
# Slide 9 — Resultados / conclusión
# ---------------------------------------------------------------------------
slide = new_slide()
add_header(slide, "RESULTADOS (100 términos, llama3.2:3b)", "Conclusión: autonomía agéntica no gana", 9)
add_table(
    slide,
    ["Método", "MRR"],
    [
        ("llm_zero_shot", "0.201"),
        ("hybrid_fusion", "0.195"),
        ("llm_rerank", "0.182"),
        ("mas_base / mas_llm_ranker / mas_safe_hybrid", "≈ 0.119 – 0.121"),
        ("supervisor (tras el fix)", "0.1215"),
    ],
    [Emu(6800000), Emu(4355680)],
    height=Emu(2400000),
)
add_bullets(slide, [
    (1, "Los métodos de una sola llamada LLM superan a todas las arquitecturas multiagente", False),
    (1, "Patrón repetido en todas las iteraciones de la tesis: más 'autonomía' agéntica no mejora los resultados frente a pipelines más simples", False),
], top=Emu(4050000), height=Emu(1800000))

# ---------------------------------------------------------------------------
# Slide 10 — Cierre
# ---------------------------------------------------------------------------
slide = new_slide()
set_bg_white(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()
bar.shadow.inherit = False

box = slide.shapes.add_textbox(MARGIN, Emu(2900000), Emu(11185855), Emu(700000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Preguntas"
r.font.size = Pt(40)
r.font.bold = True
r.font.name = "Calibri"
r.font.color.rgb = WHITE

box = slide.shapes.add_textbox(MARGIN, Emu(3700000), Emu(11185855), Emu(500000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Código y resultados: proyecto synonic/"
r.font.size = Pt(16)
r.font.name = "Calibri"
r.font.color.rgb = BLUE

out_path = "presentation/synonic_metodos.pptx"
prs.save(out_path)
print("saved:", out_path)
