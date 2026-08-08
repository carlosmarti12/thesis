"""
Generator for synonic_proyecto_completo_en.pptx - English translation of the
full project overview deck (build_proyecto_deck.py): objective, why this
rebuild exists, the leakage fix, dataset, architecture, the 11 methods,
metrics, and results from the real 100-term comparison run
(results/comparison_results.csv). Same visual system as the other decks in
this folder (dark navy header bar, Calibri, same color tokens).

Run with: .venv/bin/python presentation/build_proyecto_deck_en.py
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x11, 0x1D, 0x3A)
BLUE = RGBColor(0x3D, 0x8B, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xC9, 0xD3, 0xE0)

SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)
MARGIN = Emu(502920)
CONTENT_TOP = Emu(1234440)
CONTENT_W = SLIDE_W - 2 * MARGIN

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def set_bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_header(slide, kicker, title, page_num):
    set_bg_white(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(1051560))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    kicker_box = slide.shapes.add_textbox(MARGIN, Emu(109728), Emu(10058400), Emu(320040))
    p = kicker_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = kicker
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(MARGIN, Emu(384048), Emu(10972800), Emu(640080))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = WHITE

    foot_l = slide.shapes.add_textbox(MARGIN, Emu(6510528), Emu(5486400), Emu(274320))
    p = foot_l.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "synonic — project overview"
    r.font.size = Pt(10)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY

    foot_r = slide.shapes.add_textbox(Emu(11155680), Emu(6510528), Emu(548640), Emu(274320))
    p = foot_r.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(page_num)
    r.font.size = Pt(10)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY


def add_bullets(slide, items, top=Emu(1280160), height=Emu(5120640), left=Emu(640080), width=Emu(10881360), lvl0_size=17, lvl1_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for level, text, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        prefix = "•  " if level == 0 else "‒  "
        r = p.add_run()
        r.text = prefix + text
        r.font.name = "Calibri"
        r.font.bold = bold
        r.font.color.rgb = INK
        if level == 0:
            r.font.size = Pt(lvl0_size)
            p.space_after = Pt(8)
        else:
            r.font.size = Pt(lvl1_size)
            p.space_after = Pt(4)


def add_table(slide, headers, rows, col_widths, top=CONTENT_TOP, left=MARGIN, height=Emu(4846320), font_pt=13, header_font_pt=14, center_cols=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = sum(col_widths)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    center_cols = center_cols or set()

    for c, w in enumerate(col_widths):
        table.columns[c].width = w

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = htext
        r.font.size = Pt(header_font_pt)
        r.font.bold = True
        r.font.name = "Calibri"
        r.font.color.rgb = WHITE

    for ridx, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ridx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if c in center_cols:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = Pt(font_pt)
            r.font.bold = False
            r.font.name = "Calibri"
            r.font.color.rgb = INK


def title_slide(kicker_text, main_title, subtitle, footnote):
    slide = new_slide()
    set_bg_white(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = slide.shapes.add_textbox(MARGIN, Emu(2100000), Emu(11185855), Emu(700000))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = main_title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = WHITE

    box = slide.shapes.add_textbox(MARGIN, Emu(2870000), Emu(11185855), Emu(500000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(19)
    r.font.name = "Calibri"
    r.font.color.rgb = BLUE

    box = slide.shapes.add_textbox(MARGIN, Emu(3420000), Emu(11185855), Emu(500000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = footnote
    r.font.size = Pt(14)
    r.font.name = "Calibri"
    r.font.color.rgb = LIGHT

    box = slide.shapes.add_textbox(MARGIN, Emu(6300000), Emu(9000000), Emu(400000))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Carlos Martínez  ·  synonic/ project"
    r.font.size = Pt(12)
    r.font.name = "Calibri"
    r.font.color.rgb = GRAY
    return slide


# ===========================================================================
# Load real results
# ===========================================================================
df = pd.read_csv("results/comparison_results.csv")
agg = df.groupby("method").agg(
    exact_top1=("exact_top1", "mean"),
    exact_any=("exact_any", "mean"),
    recall_at_5=("recall_at_5", "mean"),
    mrr=("mrr", "mean"),
    avg_time_s=("time_seconds", "mean"),
    avg_candidates=("num_final_candidates", "mean"),
)
ORDER = ["same_term", "wordnet_direct", "embedding_wordnet", "llm_zero_shot",
          "llm_expansion", "llm_rerank", "hybrid_fusion", "mas_base",
          "mas_llm_ranker", "mas_safe_hybrid", "supervisor"]
agg = agg.loc[ORDER]

# ===========================================================================
# Slide 1 — Title
# ===========================================================================
title_slide(
    "SYNONIC",
    "SYNONIC — PROJECT OVERVIEW",
    "Objective, architecture, methods, and metrics for an open synonym discovery task",
    "What we did, how we did it, and how we measured it",
)

# ===========================================================================
# Slide 2 — Agenda
# ===========================================================================
slide = new_slide()
add_header(slide, "INDEX", "Agenda", 2)
add_bullets(slide, [
    (0, "Objective and task definition", True),
    (0, "Why this rebuild exists — the data leakage problem", True),
    (0, "Dataset and evaluation sample", True),
    (0, "Overall project architecture", True),
    (0, "The 11 implemented methods", True),
    (0, "How we measure — the metrics", True),
    (0, "Results (100 terms, llama3.2:3b)", True),
    (0, "Key findings and limitations", True),
    (0, "Next steps", True),
])

# ===========================================================================
# Slide 3 — Objective
# ===========================================================================
slide = new_slide()
add_header(slide, "OBJECTIVE", "Task definition", 3)
add_bullets(slide, [
    (0, "Open synonym discovery", True),
    (1, "Given an academic term (e.g. \"stocks\"), find its synonyms or alternative labels as they would appear in an academic taxonomy/database", False),
    (1, "The ground truth is the en_synonym column of the TU-Expert-Collection-Topic-Synonyms dataset", False),
    (0, "Core rule: open discovery", True),
    (1, "A method may only ever see the input term (en) — never en_synonym until the evaluation step", False),
    (1, "en_synonym is only ever used to score an output that has already been produced, never to help produce it", False),
    (0, "Why it matters", True),
    (1, "Simulates a real deployment: in production there is no \"correct answer\" to search inside", False),
])

# ===========================================================================
# Slide 4 — Why this rebuild exists (leakage)
# ===========================================================================
slide = new_slide()
add_header(slide, "METHODOLOGICAL RIGOR", "The data leakage problem and its fix", 4)
add_bullets(slide, [
    (0, "Finding in the prior iteration (new/src2/)", True),
    (1, "Several methods built their search index directly from df['en_synonym'] — the answer column", False),
    (1, "Affected: embedding retrieval, LLM expansion, and all 3 MAS variants (their retrieval_agent)", False),
    (1, "Result: the system searched literally inside the correct answer → inflated metrics, not representative of real-world use", False),
    (0, "Fix applied in synonic/", True),
    (1, "The candidate vocabulary now comes from WordNet (src/open_vocab.py) — a resource fully independent of the dataset", False),
    (1, "en_synonym is read in exactly one place in the whole codebase: run_experiment.py, inside the main loop, always after candidates are generated", False),
    (1, "Verifiable at any time: grep en_synonym src/ — every other hit is just a comment/docstring", False),
])

# ===========================================================================
# Slide 5 — Dataset
# ===========================================================================
slide = new_slide()
add_header(slide, "DATA", "Dataset and evaluation sample", 5)
add_bullets(slide, [
    (0, "Source", True),
    (1, "jensjorisdecorte/TU-Expert-Collection-Topic-Synonyms (HuggingFace)", False),
    (1, "Only topic, en (input term), and en_synonym (ground truth) are kept — nl/nl_synonym are dropped (this project targets English only)", False),
    (0, "Preparation", True),
    (1, "data/prepare_dataset.py → data/topic_synonyms_clean.csv — 970 rows after removing empty/duplicate entries", False),
    (1, "data/make_eval_sample.py → data/eval_sample.csv — a fixed, seeded 100-row sample (random_state=42)", False),
    (1, "The fixed seed guarantees every method is compared on exactly the same rows, reproducibly", False),
])

# ===========================================================================
# Slide 6 — Overall architecture
# ===========================================================================
slide = new_slide()
add_header(slide, "ARCHITECTURE", "How the project is organized", 6)
add_bullets(slide, [
    (0, "Overall flow", True),
    (1, "data/ → cleaned dataset + fixed evaluation sample", False),
    (1, "src/open_vocab.py → WordNet vocabulary index (~147k phrases), cached in data/cache/", False),
    (1, "src/llm.py → single access point to Ollama (ask_llm, parse_candidates)", False),
    (1, "src/baselines.py + src/agents.py/graph.py + src/supervisor/ → the 11 methods", False),
    (1, "src/evaluation.py → all metrics, documented", False),
    (1, "src/run_experiment.py (one method) / src/run_comparison.py (all 11 at once) → results/*.csv", False),
    (0, "Key design decision", True),
    (1, "Every method receives only term + shared resources (embedder/LLM/index) — never the full dataframe", False),
    (1, "Both CLIs are resumable: if the CSV already has rows for (term, method), they're skipped — safe to interrupt and resume", False),
])

# ===========================================================================
# Slide 7 — The 11 methods (summary table)
# ===========================================================================
slide = new_slide()
add_header(slide, "METHODS", "The 11 implemented methods (1/2)", 7)
add_table(
    slide,
    ["Method", "Description"],
    [
        ("same_term", "Returns the input term itself. A reference floor."),
        ("wordnet_direct", "Lemmas from WordNet synsets containing the term. Symbolic, no ML."),
        ("embedding_wordnet", "Embeds the term (MiniLM) + top-5 nearest neighbors by cosine similarity in the WordNet index."),
        ("llm_zero_shot", "Asks the LLM directly for 5 synonyms, no retrieval."),
        ("llm_expansion", "The LLM generates 6 query variants; each is searched in WordNet and aggregated by score."),
        ("llm_rerank", "The LLM generates 10 free candidates; they're reordered by embedding similarity."),
        ("hybrid_fusion", "3 LLM calls (self-consistency) + WordNet retrieval, fused via weighted RRF."),
    ],
    [Emu(2800000), Emu(8355680)],
    height=Emu(4600000),
    font_pt=13,
)

# ===========================================================================
# Slide 8 — The 11 methods (MAS + supervisor)
# ===========================================================================
slide = new_slide()
add_header(slide, "METHODS", "The 11 implemented methods (2/2) — MAS and supervisor", 8)
add_bullets(slide, [
    (0, "mas_base (v1)", True),
    (1, "generator(LLM) + retrieval(WordNet) → merge → verification (0.75·semantic + 0.25·fuzzy) → top 5", False),
    (0, "mas_llm_ranker (v2)", True),
    (1, "retrieval → generator → merge → an LLM reorders the whole list → top 5. Confirmed negative ablation", False),
    (0, "mas_safe_hybrid (v3)", True),
    (1, "Same as v2 + safe_finalizer: 0.9·embedding_sim + 0.1·LLM bonus (never penalizes)", False),
    (0, "supervisor", True),
    (1, "A single LLM decides, turn by turn, which action to take (retrieve/generate/merge/score/finalize/finish)", False),
    (1, "Reuses the same deterministic functions as the MAS variants — only who decides the order changes", False),
    (1, "A real planning bug was found and fixed using the project's own comparison harness (see findings)", False),
])

# ===========================================================================
# Slide 9 — Metrics
# ===========================================================================
slide = new_slide()
add_header(slide, "HOW WE MEASURE", "Metric definitions", 9)
add_table(
    slide,
    ["Metric", "What it measures"],
    [
        ("exact_top1", "Candidate #1 matches the ground truth exactly."),
        ("exact_any / recall@5", "The ground truth appears anywhere among the 5 returned candidates."),
        ("recall_at_3", "The ground truth appears among the top 3 candidates."),
        ("precision_at_5", "Only one ground-truth synonym exists per term → mechanically capped at 0.2; equivalent to recall@5 / 5."),
        ("mrr", "1 / rank of the correct answer (0 if absent). Rewards ranking it near the top even if not #1."),
        ("fuzzy_similarity", "rapidfuzz.fuzz.ratio — lexical/spelling closeness (e.g. \"optimise\" vs \"optimize\")."),
        ("semantic_similarity", "Cosine similarity of embeddings (MiniLM) — meaning-based closeness, not textual form."),
    ],
    [Emu(2900000), Emu(8255680)],
    height=Emu(4600000),
    font_pt=13,
)

# ===========================================================================
# Slide 10 — Results (full table, real data)
# ===========================================================================
slide = new_slide()
add_header(slide, "RESULTS", "Full comparison — 100 terms, llama3.2:3b", 10)

rows = []
for m in ORDER:
    row = agg.loc[m]
    rows.append((
        m,
        f"{row['mrr']:.3f}",
        f"{row['exact_any']*100:.0f}%",
        f"{row['recall_at_5']*100:.0f}%",
        f"{row['avg_time_s']:.2f}s",
    ))

add_table(
    slide,
    ["Method", "MRR", "exact_any", "recall@5", "Time/term"],
    rows,
    [Emu(3200000), Emu(1996464), Emu(1996464), Emu(1996464), Emu(1996464)],
    height=Emu(4600000),
    font_pt=13,
    header_font_pt=13,
    center_cols={1, 2, 3, 4},
)

# ===========================================================================
# Slide 11 — Key findings
# ===========================================================================
slide = new_slide()
add_header(slide, "FINDINGS", "Reading the results", 11)
add_bullets(slide, [
    (0, "The pattern repeats across every iteration of the thesis", True),
    (1, "Single-LLM-call methods (llm_zero_shot, hybrid_fusion, llm_rerank) beat full multi-agent architectures on MRR", False),
    (1, "More agentic \"autonomy\" (MAS, supervisor) doesn't improve results over simpler pipelines — and costs more time", False),
    (0, "hybrid_fusion is the best on recall@5 / exact_any (31%)", True),
    (1, "Combining self-consistency LLM generation with WordNet retrieval via RRF wins on coverage, though not on pure MRR", False),
    (0, "mas_llm_ranker (v2) confirms the negative ablation", True),
    (1, "Letting the LLM freely reorder the merged list makes results worse than a numeric ranking (v1) or hybrid one (v3)", False),
    (0, "supervisor: a real planning bug, found and fixed", True),
    (1, "The LLM sometimes generated candidates after merging, and those candidates were left orphaned, never re-merged", False),
    (1, "After the fix (checking order, not just presence, of actions): mrr 0.091 → 0.122, recall@5 0.13 → 0.21", False),
    (0, "Time cost", True),
    (1, "From 0s (same_term) to 7.2s/term (supervisor) — the slowest method is not the best-performing one", False),
])

# ===========================================================================
# Slide 12 — Limitations
# ===========================================================================
slide = new_slide()
add_header(slide, "LIMITATIONS", "What this analysis doesn't cover (yet)", 12)
add_bullets(slide, [
    (0, "Small model", True),
    (1, "llama3.2:3b via Ollama — a larger model would likely improve every LLM-dependent method with no architectural change", False),
    (0, "WordNet coverage ceiling", True),
    (1, "It doesn't contain every domain-specific term the dataset expects — this penalizes index-anchored methods (retrieval, MAS) versus purely generative ones", False),
    (0, "Idiosyncratic ground truths", True),
    (1, "Some reference answers are highly specific editorial choices that no method can predict from the term alone", False),
    (0, "Sample scale", True),
    (1, "100 of 970 available terms — validating at larger scale is the natural next step", False),
])

# ===========================================================================
# Slide 13 — Next steps
# ===========================================================================
slide = new_slide()
add_header(slide, "CLOSING", "Next steps", 13)
add_bullets(slide, [
    (0, "Validate at larger scale", True),
    (1, "Rerun the full comparison over all 970 rows of the dataset, not just the 100-row sample", False),
    (0, "Try a larger model", True),
    (1, "Measure how much of the gap between methods is due to LLM size rather than architecture", False),
    (0, "Go deeper on supervisor", True),
    (1, "Assess whether splitting the supervisor into independent cognitive agents (disambiguation / generation / critique) adds real value before building it", False),
    (0, "Document for the thesis", True),
    (1, "Carry this comparison and its findings into the results chapter, with the full tables from comparison_results.csv", False),
])

# ===========================================================================
# Slide 14 — Closing
# ===========================================================================
slide = new_slide()
set_bg_white(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()
bar.shadow.inherit = False

box = slide.shapes.add_textbox(MARGIN, Emu(2900000), Emu(11185855), Emu(700000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Questions"
r.font.size = Pt(40)
r.font.bold = True
r.font.name = "Calibri"
r.font.color.rgb = WHITE

box = slide.shapes.add_textbox(MARGIN, Emu(3700000), Emu(11185855), Emu(500000))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Code and results: synonic/ project"
r.font.size = Pt(16)
r.font.name = "Calibri"
r.font.color.rgb = BLUE

out_path = "presentation/synonic_proyecto_completo_en.pptx"
prs.save(out_path)
print("saved:", out_path)
print(agg.to_string())
